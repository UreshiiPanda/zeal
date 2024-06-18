#!/usr/bin/env python
# coding: utf-8

# In[1]:


# pip install packages
get_ipython().run_line_magic('pip', 'uninstall openai --yes')
get_ipython().run_line_magic('pip', 'install openai==0.28')
get_ipython().run_line_magic('pip', 'install tiktoken')
get_ipython().run_line_magic('pip', 'install pdfminer')
get_ipython().run_line_magic('pip', 'install pdfminer.six')
get_ipython().run_line_magic('pip', 'install docx2txt')
get_ipython().run_line_magic('pip', 'install python-pptx')
# %pip install pinecone-client==2.2.4
get_ipython().run_line_magic('pip', 'install pinecone')
get_ipython().run_line_magic('pip', 'install pinecone-client --upgrade')



# imports
import openai
#from openai import OpenAI
import os
import tiktoken
from tqdm.auto import tqdm     # this is our progress bar
import pinecone
import re
import pandas as pd
from pdfminer.high_level import extract_text, extract_pages
import docx2txt
from pptx import Presentation
import zipfile
import xml.etree.ElementTree as ET
import base64



# constants
GPT_MODEL = 'gpt-4o'
EMBEDDING_MODEL = 'text-embedding-ada-002'

# this is for the input data
PATH = './data'

# these are Pinecone vars
INDEX = 'idx'


###### general info
# the purpose of calculating the number of tokens is to ensure that the text chunks fit within the token limit of the model being used
# The notebook splits the text into chunks based on a token limit of 256. 
    # so this is what decides how long each string in the res arr will be
    # The choice of 256 tokens as the chunk size could be based on a balance between granularity and efficiency
    # Smaller chunk sizes would result in more chunks and more embeddings to store, which could increase storage requirements and processing time
    # Larger chunk sizes might lead to less specific embeddings and potentially miss important details within the text.
# The size of each line depends on the original structure of the text files being processed.
# The text-embedding-ada-002 model produces embeddings with a dimension of 1536, so the Pinecone index is created with the same dimension to ensure compatibility.
# The batch size is set to 128 for generating embeddings using the OpenAI API.
    # The batch size is set to 32 when uploading the embeddings and metadata to Pinecone using the upsert operation.
    # Larger batch sizes can be more efficient in terms of API calls and processing time, but they also require more 
    # memory and may be limited by the API's constraints
# The max_tokens parameter determines the maximum number of tokens to retrieve from the relevant chunks based on the input query.
    # Retrieving more tokens provides more context but also increases the response size and processing time


# In[2]:


get_ipython().run_line_magic('pip', 'install python-dotenv')

from dotenv import load_dotenv

load_dotenv()

ENV = os.getenv("ENV")
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
PINECONE_API_KEY = os.getenv('PINECONE_API_KEY')

openai.api_key = os.getenv('OPENAI_API_KEY')


# In[3]:


def num_tokens(text: str, model: str = 'gpt-4') -> int:
    encoding = tiktoken.encoding_for_model(model)
    return len(encoding.encode(text))


# In[4]:


def read_txt(path):
    with open(path, encoding='utf8') as f:
        return f.readlines()


# In[5]:


def read_pdf(path):
    return extract_text(path).split('\n')


# In[6]:


def read_docx(path):
    # extract text
    return docx2txt.process(path).split('\n')


# In[7]:


def read_pptx(path):
    split = []
    prs = Presentation(path)
    print("----------------------")
    temp = ''
    enterCount = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
#                 if len(shape.text) > 0:
#                     temp += shape.text
#                 elif enterCount > 1:
#                     split.append(temp)
#                     temp = ''
#                     enterCount = 0
#                 elif len(shape.text) == 0:
#                     enterCount += 1
                split.append(shape.text)
    
    return split


# In[8]:


def read_xlsx(path):
    with zipfile.ZipFile(path, 'r') as zip_ref:
        zip_ref.extractall(PATH)
    
    xml_file = os.path.join(PATH, 'xl/sharedStrings.xml')
    
    tree = ET.parse(xml_file)
    root = tree.getroot()
    
    result = []
    
    for elem in root.iter():
        if elem.text:
            result.append(elem.text.strip())
    
    return result


# In[9]:


def read_png(path):
    with open(path, "rb") as image_file:
        image = base64.b64encode(image_file.read()).decode("utf-8")
        
        openai.api_key = os.environ['OPENAI_API_KEY']
        
        messages = [
            { 'role': 'system', 'content': 'system message here' },
            { 'role': 'user', 'content': [
                {'type': 'text', 'text': f'Convert the following image into a text table.'},
                {'type': 'image_url', 'image_url': { 'url': f'data:image/png;base64,{image}'} }
            ]}
        ]
        
        functions = [
            {
                'name': 'write_queries',
                'description': 'Get information from the database',
                'parameters': {
                    'type': 'object',
                    'properties': {
                        'details': {
                            'type': 'string',
                            'description': 'The information needed from the database',
                        },
                    }
                }
            }
        ]
        
        result = ''
        
        for chunk in openai.chat.completions.create(
            model=GPT_MODEL,
            messages=messages,
            functions=functions,
            temperature=0,
            frequency_penalty=0,
            presence_penalty=0,
            function_call=None,
            stream=True
        ):
            if len(chunk.choices) > 0:
                response = chunk.choices[0]
                finish_reason = response.finish_reason
                content = response.delta.content
                function = response.delta.function_call

                if content is not None:
                    print(content, end='')
                    result += content
        
        return result.split('\n')


# In[10]:


def collect_all_files(root_path):
    all_files = []
    for root, dirs, files in os.walk(root_path):
        for file in files:
            file_path = os.path.join(root, file)
            all_files.append(file_path)
    return all_files


# In[11]:


def get_range(s, delimiter = '.'):
    parts = s.split(delimiter)
    return delimiter.join(parts[-1:])


# In[12]:


messages = []

files = collect_all_files(PATH)

lines = []
for file in files:
    lines = []
    print(file)
    suffix = get_range(file)
    if suffix == 'txt':
        lines = read_txt(file)
    elif suffix == 'pdf':
        lines = read_pdf(file)
    elif suffix == 'docx':
        lines = read_docx(file)
    elif suffix == 'pptx':
        lines = read_pptx(file)
    elif suffix == 'xlsx':
        lines = read_xlsx(file)
    elif suffix == 'png':
        lines = read_png(file)

    temp_string = []
    # print(lines)
    for line in lines:
        if num_tokens('\n'.join(temp_string)) > 256:
            # print('\n'.join(temp_string))
            messages.append('\n'.join(temp_string))
            temp_string = []
        else:
            temp_string.append(line)
    if num_tokens('\n'.join(temp_string)) < 256:
        # print('\n'.join(temp_string))
        messages.append('\n'.join(temp_string))


# In[13]:


messages = [message for message in messages if len(message) > 0]


# In[14]:


if not messages: 
    # test out messages
    messages = ["this is message 1", "this is message 2", "this is message 3", "this is message 4", "meow meow", "woof woof", "I am a cat"]

# print(messages)
# for i in range(len(messages)):
#     print(len(messages[i]))
#     print(messages[i])
#     print(" \n###########################\n ")


# In[15]:


from pinecone import Pinecone, PodSpec

INDEX = 'idx'

# pc = Pinecone(api_key=os.environ['PINECONE_API_KEY'])
pc = Pinecone(api_key=os.getenv('PINECONE_API_KEY'))
# print(f"{os.environ['PINECONE_API_KEY']}")

if INDEX not in pc.list_indexes().names():
    pc.create_index(INDEX, dimension=1536, metric='cosine', spec=PodSpec(environment=ENV))
# connect to index
index = pc.Index(INDEX)
print(index)


# In[16]:


# calculate embeddings
BATCH_SIZE = 128  # you can submit up to 2048 embedding inputs per request

embeddings = []
for batch_start in range(0, len(messages), BATCH_SIZE):
    batch_end = batch_start + BATCH_SIZE
    batch = messages[batch_start:batch_end]
    print(f"Batch {batch_start} to {batch_end-1}")
    response = openai.Embedding.create(model=EMBEDDING_MODEL, input=batch)
    for i, be in enumerate(response['data']):
        assert i == be['index']  # double check embeddings are in same order as input
    batch_embeddings = [e['embedding'] for e in response['data']]
    embeddings.extend(batch_embeddings)

df = pd.DataFrame({'text': messages, 'embedding': embeddings})


# In[17]:


print(index)
batch_size = 32  # process everything in batches of 32
for i in tqdm(range(0, len(df['text']), batch_size)):
    # set end position of batch
    i_end = min(i+batch_size, len(df['text']))
    # get batch of lines and IDs
    lines_batch = df['text'][i:i+batch_size]
    embeds = df['embedding'][i:i+batch_size]
    ids_batch = ['training:' + str(n) for n in range(i, i_end)]
    # prep metadata and upsert batch
    meta = [{'text': line} for line in lines_batch]
    to_upsert = zip(ids_batch, embeds, meta)
    # upsert to Pinecone
    index.upsert(vectors=list(to_upsert))


# In[18]:


embeddings = df.to_dict()


# In[19]:


def send_embeddings(query: str = '', embeddings: dict = {}, max_tokens: int = 1024) -> str:
    '''Return the max_tokens amount of related contexts based on the query string.

    Keyword arguments:
    query -- The query needing context
    embeddings -- The list of embeddings and text
    max_tokens -- The number of tokens to pull (Default 1024)

    Returns:
    The context of the query
    '''
    embedding = openai.Embedding.create(model=EMBEDDING_MODEL, input=query)['data'][0]['embedding']
    context = [message for _, message in sorted(zip(cosine_similarity(embeddings['embedding'], [embedding]), embeddings['text']), reverse=True)]


    max_tokens = int(max_tokens)

    text = ''
    total_tokens = 0
    max_tokens = min(num_tokens('\n'.join(context)) + total_tokens + 1, max_tokens + 1)
    i = 0

    for item in context:
        nxt = f'\n{item}\n'
        total_tokens += num_tokens(nxt)
        if total_tokens > max_tokens:
            break
        text += nxt
        i += 1

    return f'{text}'


# In[20]:


# chatbox = ''
# while chatbox != 'exit':
#     chatbox = input()
    
#     if chatbox == 'exit':
#         continue
    
#     index = "idx"

#     # client = OpenAI(api_key=os.environ['OPENAI_API_KEY'])

#     # STEP 1: Embed your prompt   (create an embedding from your prompt via OpenAI's embeddings)
#     embedding = openai.Embedding.create(model=EMBEDDING_MODEL, input=chatbox).data[0].embedding

#     # STEP 2: Initialize pinecone index
#     pc = Pinecone(api_key=os.environ['PINECONE_API_KEY'])

#     if index in pc.list_indexes().names():
#         idx = pc.Index(index)
#         result = idx.query(vector=[embedding], top_k=3, include_metadata=True)
#         # replace() added for the newlines
#         context = [x['metadata']['text'].replace('\n', '') for x in result['matches']]


#         print(context) 


# In[21]:


chatbox = ''
while chatbox != 'exit':
    chatbox = input()    
    if chatbox == 'exit':
        continue    
    index = "idx"    
    
    # if using newer version of openAI import:    client = OpenAI(api_key=os.environ['OPENAI_API_KEY'])    
    # STEP 1: Embed your prompt
    embedding = openai.Embedding.create(model=EMBEDDING_MODEL, input=chatbox).data[0].embedding    
    
    # STEP 2: Initialize pinecone index
    pc = Pinecone(api_key=os.environ['PINECONE_API_KEY'])    


    ### increase top_k to increase the num of vectors used from PineCone (so top_k=3 would only use the top 3 similar vectors
    if index in pc.list_indexes().names():
        idx = pc.Index(index)
        result = idx.query(vector=[embedding], top_k=1024, include_metadata=True)
        context = [x['metadata']['text'].replace('\n', '') for x in result['matches']]
        #print(context)


    # STEP 3: Generate a response using the following pre-exisiting context from Pinecone
    messages = [
        {"role": "system", "content": "You are a helpful assistant."},
        {"role": "user", "content": chatbox},
        {"role": "assistant", "content": "Based on the following context:"},
        {"role": "system", "content": " ".join(context)}
    ]    
    
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=messages,
        # more tokens means longer/more-detailed responses and VV
        max_tokens=1500,
        # more temp is a more subjective/creative/random response, less temp is more objective/direct/deterministic response
        temperature=0.7
    )    
    
    # Extract and print the response text
    response_text = response.choices[0].message['content'].strip()
    print(f"\n\nWelcome to Costco, I love you:\n\n {response_text}")

